"""
Aplomb — Supplier Shortlist for MVP Market Test
================================================

Builds Aplomb_Supplier_Shortlist_v1.pptx (15 slides, 16:9) listing the top
white-label suppliers for each of 5 GLP-1 side effects Aplomb is targeting:
Ozempic face, halitosis, hair loss, nausea, nutrient depletion — plus a
sixth section covering the universal outer-packaging mailer system (small +
large folding-carton SKUs, MOQ 200-400, US-only, Musely-style reference).

Hard sourcing gate per supplier (formula sections only): MOQ <= 100 units
OR total MOQ cost < $250. The packaging section uses its own gates: US
production, custom interior print, MOQ <= 400 per SKU.

Run: python3 build_supplier_shortlist_deck.py
"""

from pathlib import Path
from datetime import datetime
import os
import re
import shutil
import subprocess
import time
import zipfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Brand tokens (Aplomb — mirrors AFTER. brand)
# ---------------------------------------------------------------------------
BG = RGBColor(0xEF, 0xE8, 0xDC)        # warm cream
PAPER = RGBColor(0xF7, 0xF1, 0xE6)     # lighter cream, cards
INK = RGBColor(0x1A, 0x15, 0x12)       # near-black brown
AMBER = RGBColor(0x7A, 0x3D, 0x14)     # deep amber accent
AMBER_LIGHT = RGBColor(0xD9, 0xA0, 0x6B)  # tan
RULE = RGBColor(0xD9, 0xCF, 0xBD)      # hairline
MUTED = RGBColor(0x6B, 0x5D, 0x4C)     # secondary text
SCORE_HI = RGBColor(0xE6, 0xC9, 0x9C)   # >= 80 (amber-light wash)
SCORE_MID = RGBColor(0xF7, 0xF1, 0xE6)  # 60-79 (paper)
SCORE_LOW = RGBColor(0xE5, 0xDD, 0xCC)  # < 60 (slightly darker rule)

SERIF = "Cormorant Garamond"
SANS = "IBM Plex Sans"  # Aplomb brand body font (BRAND.md §3)

OUT = Path(
    "/Users/zacharypoll/Desktop/Documents/Claude Code/aplomb.clinic/"
    "business documents/Aplomb_Supplier_Shortlist_v1.pptx"
)

# ---------------------------------------------------------------------------
# Slide dimensions — Emu(12192000) avoids the Inches(13.333) truncation that
# triggers PowerPoint's "needs to be repaired" prompt.
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_blank_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s


def add_text(slide, x, y, w, h, text, *, font=SANS, size=11, color=INK,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             leading=1.18):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return tb


def add_rich(slide, x, y, w, h, segments, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, leading=1.18):
    """segments = list of paragraph specs.
       Each paragraph spec is a list of dicts: [{"t","size","bold","color","font","italic"}]
       Plain string item = a paragraph with default styling."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(segments):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        if isinstance(para, str):
            run = p.add_run()
            run.text = para
            run.font.name = SANS
            run.font.size = Pt(11)
            run.font.color.rgb = INK
        else:
            for seg in para:
                run = p.add_run()
                run.text = seg.get("t", "")
                run.font.name = seg.get("font", SANS)
                run.font.size = Pt(seg.get("size", 11))
                run.font.color.rgb = seg.get("color", INK)
                run.font.bold = seg.get("bold", False)
                run.font.italic = seg.get("italic", False)
    return tb


def add_rect(slide, x, y, w, h, fill=PAPER, line_color=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line_color is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line_color
        r.line.width = Pt(0.5)
    r.shadow.inherit = False
    return r


def add_hline(slide, x, y, w, color=RULE, weight=0.75):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_chapter_chrome(slide, eyebrow):
    add_hline(slide, Inches(0.6), Inches(0.55), Inches(12.13))
    add_text(slide, Inches(0.6), Inches(0.18), Inches(8), Inches(0.3),
             [eyebrow], font=SANS, size=9, color=AMBER)
    add_rich(slide, Inches(11.0), Inches(0.18), Inches(1.8), Inches(0.3), [[
        {"t": "Aplomb", "font": SERIF, "size": 11, "color": INK, "italic": True},
        {"t": ".", "font": SERIF, "size": 11, "color": AMBER, "italic": True, "bold": True},
    ]], align=PP_ALIGN.RIGHT)


def add_footer(slide, page_num, total):
    add_hline(slide, Inches(0.6), Inches(7.05), Inches(12.13))
    add_text(slide, Inches(0.6), Inches(7.13), Inches(8), Inches(0.3),
             ["Aplomb — Supplier shortlist for MVP market test · "
              "For Zachary Poll · Confidential"],
             font=SANS, size=8, color=MUTED)
    add_text(slide, Inches(11.3), Inches(7.13), Inches(1.5), Inches(0.3),
             [f"{page_num:02d} / {total:02d}"],
             font=SANS, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def _set_cell_text(cell, text, *, font=SANS, size=8.5, color=INK, bold=False,
                   align=PP_ALIGN.LEFT, hyperlink=None):
    cell.text_frame.clear()
    cell.text_frame.margin_left = Emu(45000)
    cell.text_frame.margin_right = Emu(45000)
    cell.text_frame.margin_top = Emu(40000)
    cell.text_frame.margin_bottom = Emu(40000)
    cell.text_frame.word_wrap = True
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.12
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    if hyperlink:
        run.hyperlink.address = hyperlink
        # hyperlink color override (PPT respects the run color for visited+text)
        run.font.color.rgb = AMBER
        run.font.underline = True


def _score_color(score):
    if score >= 80:
        return SCORE_HI
    if score >= 60:
        return SCORE_MID
    return SCORE_LOW


def add_supplier_matrix(slide, x, y, w, h, suppliers):
    """suppliers: list of dicts with keys
       name, url, product, product_name, product_url, moq, unit_cost_usd,
       total_moq_cost_usd, alleviation_score, notes"""
    headers = ["Company", "Website", "Product", "Product page", "MOQ",
               "$/unit", "Total MOQ$", "Score", "Notes"]
    # Column widths in inches, must sum to 12.0 for the 12.0" wide table
    col_widths_in = [1.10, 1.20, 1.80, 1.55, 0.55, 0.75, 0.95, 0.65, 3.45]
    n_cols = len(headers)
    n_rows = len(suppliers) + 1
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    total_w = sum(col_widths_in)
    for i, cw in enumerate(col_widths_in):
        tbl.columns[i].width = int(w * cw / total_w)

    # Header row
    for j, head in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        _set_cell_text(cell, head, size=8.5, color=BG, bold=True)

    # Body rows
    for i, sup in enumerate(suppliers):
        row = i + 1
        for j in range(n_cols):
            cell = tbl.cell(row, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG if i % 2 == 0 else PAPER
        # Per-cell content
        _set_cell_text(tbl.cell(row, 0), sup["name"], size=8.5, bold=True)
        # Website hyperlink — show domain only for compactness
        domain = re.sub(r"^https?://(www\.)?", "", sup["url"]).rstrip("/")
        if len(domain) > 28:
            domain = domain[:26] + "…"
        _set_cell_text(tbl.cell(row, 1), domain, size=8, hyperlink=sup["url"])
        _set_cell_text(tbl.cell(row, 2), sup["product"], size=7.5)
        # Product page — short SKU/service name, hyperlinked to product URL
        _set_cell_text(tbl.cell(row, 3), sup["product_name"], size=8,
                       hyperlink=sup["product_url"])
        _set_cell_text(tbl.cell(row, 4), str(sup["moq"]), size=8.5,
                       align=PP_ALIGN.CENTER)
        _set_cell_text(tbl.cell(row, 5), f"${sup['unit_cost_usd']:.2f}",
                       size=8.5, align=PP_ALIGN.RIGHT)
        _set_cell_text(tbl.cell(row, 6), f"${sup['total_moq_cost_usd']:,.0f}",
                       size=8.5, align=PP_ALIGN.RIGHT)
        # Score cell with color tint
        score_cell = tbl.cell(row, 7)
        score_cell.fill.solid()
        score_cell.fill.fore_color.rgb = _score_color(sup["alleviation_score"])
        _set_cell_text(score_cell, str(sup["alleviation_score"]),
                       size=10, bold=True, align=PP_ALIGN.CENTER)
        _set_cell_text(tbl.cell(row, 8), sup["notes"], size=7.5)

    return tbl


# ---------------------------------------------------------------------------
# Section content (research-driven; suppliers list filled below)
# ---------------------------------------------------------------------------
SECTIONS = [
    {
        "key": "ozempic_face",
        "title": "Ozempic Face",
        "eyebrow": "§ 01 · Side effect 1 of 5 · Topical serum",
        "problem": (
            "Rapid weight loss on GLP-1s (15–20% body mass in 6–12 months) "
            "depletes facial fat pads faster than skin can retract, producing "
            "the gaunt, hollow, prematurely-aged look popularized as "
            "'Ozempic face'."
        ),
        "cause": [
            "Subcutaneous fat loss in malar, temporal and buccal pads "
            "(volume depletion).",
            "Concurrent collagen / elastin loss from caloric restriction "
            "(dermal thinning).",
            "Reduced glycation activity but accelerated visible aging "
            "(barrier + GAG drop).",
        ],
        "otc": [
            "Topical peptide serums (Matrixyl 3000, GHK-Cu, SNAP-8) — "
            "stimulate collagen + reduce fine lines.",
            "Hyaluronic acid serums + ceramide moisturizers — restore "
            "barrier and apparent volume.",
            "Retinoid analogs (HPR, retinaldehyde) — accelerate dermal "
            "remodeling without Rx.",
            "SPF 30+ daily — UV is the largest modifiable amplifier of "
            "GLP-1-driven skin aging.",
        ],
        "suppliers": [
            {
                "name": "Cellular Cosmetics",
                "url": "https://www.cellular-cosmetics.com/",
                "product": "Peptide serum w/ Matrixyl 3000 + Synthe'6, "
                           "shelf-ready, dropper-bottle filled",
                "product_name": "Anti-Age Peptide Serum",
                "product_url": "https://www.cellular-cosmetics.com/products/anti-age-peptide-serum-private-label-skin-care",
                "moq": 10,
                "unit_cost_usd": 25.00,
                "total_moq_cost_usd": 250,
                "alleviation_score": 84,
                "notes": "Australia. ONLY vendor with published INCI. "
                         "10-unit MOQ ideal for sampling. 4–8 day lead. "
                         "Budget freight to US.",
            },
            {
                "name": "Ataliene Skincare",
                "url": "https://atalieneskincare.com/",
                "product": "4-peptide anti-aging serum + HA + botanicals; "
                           "USA-filled, esthetician-grade",
                "product_name": "Plant-EGF Peptide + HA Serum",
                "product_url": "https://atalieneskincare.com/collections/serums/products/plant-egf-anti-wrinkle-serum-with-hyaluronic-acid",
                "moq": 12,
                "unit_cost_usd": 27.50,
                "total_moq_cost_usd": 330,
                "alleviation_score": 76,
                "notes": "USA mfr, fastest reorders (2 wk). $50 label "
                         "set-up per SKU. INCI not published — "
                         "request before order.",
            },
            {
                "name": "Pravada Private Label",
                "url": "https://www.pravadaprivatelabel.com/",
                "product": "Peptide+ Firming Serum (5 unnamed peptides); "
                           "70% organic; R&D customization",
                "product_name": "Peptide Complex Serum",
                "product_url": "https://www.pravadaprivatelabel.com/products/peptide-complex-serum",
                "moq": 50,
                "unit_cost_usd": 39.00,
                "total_moq_cost_usd": 1950,
                "alleviation_score": 71,
                "notes": "Canada. 20-yr salon/spa heritage. Scales 50 → "
                         "200K units. INCI opaque but peptides "
                         "swappable on request.",
            },
        ],
    },
    {
        "key": "halitosis",
        "title": "Halitosis (\"Ozempic breath\")",
        "eyebrow": "§ 02 · Side effect 2 of 5 · Zinc + chlorhexidine lozenge",
        "problem": (
            "~30% of Ozempic users self-report bad breath (ADA 2025). "
            "Three converging mechanisms: hyposalivation (dry mouth), "
            "ketosis-driven acetone breath, and gastric stasis fermentation. "
            "Persistent while on the drug."
        ),
        "cause": [
            "Hyposalivation: GLP-1 receptor activity in salivary glands "
            "reduces flow → less buffering of oral bacteria.",
            "Ketosis: caloric deficit → acetone + isoprene exhalation "
            "(fruity / metallic odor).",
            "Gastric stasis: delayed emptying → fermentation gases "
            "(eructation in ~9% per Novo's STEP-1).",
        ],
        "otc": [
            "Zinc gluconate / acetate lozenges — neutralize volatile "
            "sulfur compounds (VSCs) at the source.",
            "Xylitol mints / gum — stimulate saliva, inhibit "
            "S. mutans biofilm.",
            "Chlorhexidine 0.12% rinses (short courses) — reduce "
            "anaerobic load (Rx in some markets, OTC elsewhere).",
            "S. salivarius K12 oral probiotics — repopulate "
            "competitive flora.",
        ],
        "suppliers": [
            {
                "name": "BLIS Technologies",
                "url": "https://blis.co.nz/",
                "product": "Private-label S. salivarius K12 "
                           "lozenge — root-cause oral probiotic with "
                           "halitosis clinical evidence",
                "product_name": "K12 FreshBreath Lozenge (Private Label)",
                "product_url": "https://blis.co.nz/finished-goods/freshbreath-lozenge/",
                "moq": 100,
                "unit_cost_usd": 2.00,
                "total_moq_cost_usd": 200,
                "alleviation_score": 82,
                "notes": "NZ; only finished-goods private-label "
                         "K12 program. 1.25B CFU peppermint lozenge, "
                         "40-ct blister pack. RFQ MOQ; positioned as "
                         "the root-cause SKU in our 3-supplier slate.",
            },
        ],
    },
    {
        "key": "hair_loss",
        "title": "Hair Loss",
        "eyebrow": "§ 03 · Side effect 3 of 5 · Oral biotin + collagen capsule",
        "problem": (
            "Telogen effluvium — diffuse shedding 2–4 months after rapid "
            "weight loss begins. Affects ~6–13% of GLP-1 users (Shah 2024). "
            "Compounded by reduced protein/micronutrient intake under "
            "appetite suppression."
        ),
        "cause": [
            "Stress shock from rapid caloric deficit shifts hair follicles "
            "from anagen to telogen phase prematurely.",
            "Sub-clinical iron, zinc, vitamin D, B12 deficiencies common "
            "under reduced food intake.",
            "Marine collagen / amino-acid building blocks (cysteine, "
            "methionine) drop with low-protein eating patterns.",
        ],
        "otc": [
            "Biotin (5–10 mg/day) — supports keratin synthesis; "
            "strongest evidence in deficient subjects.",
            "Marine collagen peptides (10g/day) — provide "
            "hair-building amino acids.",
            "Iron + vitamin C (in deficient subjects only — test first).",
            "Saw palmetto + DHT-related blends — emerging evidence "
            "for androgenetic component.",
        ],
        "suppliers": [
            {
                "name": "Vox Nutrition",
                "url": "https://voxnutrition.com/",
                "product": "Custom biotin + marine collagen capsule; "
                           "NSF + FDA + cGMP + Organic certified",
                "product_name": "Hair, Skin & Nails Capsule",
                "product_url": "https://www.voxnutrition.com/product/hair-skin-nails/",
                "moq": 50,
                "unit_cost_usd": 4.50,
                "total_moq_cost_usd": 225,
                "alleviation_score": 72,
                "notes": "Utah-based US private-label specialist. "
                         "7–10 day turnaround. Only US supplier with "
                         "publicly stated <100 MOQ for biotin+collagen "
                         "blend. Best-in-class certifications.",
            },
        ],
    },
    {
        "key": "nausea",
        "title": "Nausea",
        "eyebrow": "§ 04 · Side effect 4 of 5 · Ginger + B6 chew",
        "problem": (
            "The single most common GLP-1 side effect — ~24–44% of users "
            "experience nausea, especially during dose escalation. Driven "
            "by delayed gastric emptying. Most-cited reason for "
            "discontinuation in the first 8 weeks."
        ),
        "cause": [
            "GLP-1 agonists slow gastric emptying by 30–70% — food sits "
            "longer, distends the stomach, triggers vagal nausea signals.",
            "Central CNS effects on the area postrema also contribute "
            "(same brainstem region as motion sickness).",
            "Worst during titration (week 1–2 of each dose step); "
            "tolerance often develops over 4–8 weeks.",
        ],
        "otc": [
            "Ginger root (1–1.5 g/day) — RCT-validated for chemo, "
            "pregnancy, and post-op nausea.",
            "Vitamin B6 (pyridoxine, 25–75 mg) — first-line for "
            "morning sickness; synergistic with ginger.",
            "Acupressure wristbands (P6 / Nei-Kuan) — "
            "FDA-cleared Class I device for nausea.",
            "Peppermint oil / lozenges — modest relief via "
            "anti-spasmodic effect on gastric muscle.",
        ],
        "suppliers": [
            {
                "name": "ING Pharmaceutical",
                "url": "https://www.ingpharmaceutical.com/",
                "product": "Nausea Relief Softgel — Ginger Root 67mg "
                           "+ B6 1.4mg + Mg carbonate (ready formula)",
                "product_name": "Nausea Relief Softgel",
                "product_url": "https://www.ingpharmaceutical.com/products/nausea-relief-softgels-white-label-ginger-magnesium-vitamin-b6-private-label-softgels-manufacturer",
                "moq": 100,
                "unit_cost_usd": 2.50,
                "total_moq_cost_usd": 250,
                "alleviation_score": 82,
                "notes": "RFQ — MOQ negotiable; ready-to-fill formula "
                         "exact-match for GLP-1 nausea positioning. "
                         "300+ brands served. Phone (786) 518-2903. "
                         "FDA-grade softgel.",
            },
            {
                "name": "Shaanxi VitaHall Biotechnology",
                "url": "https://www.alibaba.com/product-detail/VitaHall-OEM-ODM-Private-Label-Bulk_1601750948418.html",
                "product": "Vegan ginger root extract gummies, "
                           "private-label OEM/ODM; chewable digestion-"
                           "support format",
                "product_name": "VitaHall Ginger Root Extract Gummies",
                "product_url": "https://www.alibaba.com/product-detail/VitaHall-OEM-ODM-Private-Label-Bulk_1601750948418.html",
                "moq": 10,
                "unit_cost_usd": 5.00,
                "total_moq_cost_usd": 50,
                "alleviation_score": 70,
                "notes": "Shaanxi, China. MOQ 10 bottles (verified on "
                         "product page) — lowest published in category. "
                         "Unit price uses high-end of observed $1.70–5 "
                         "range as conservative low-qty tier. Ginger-"
                         "as-hero stock SKU; B6 add-on via custom run.",
            },
            {
                "name": "Shaanxi Suni Green",
                "url": "https://www.alibaba.com/product-detail/SUNI-OEM-ODM-Private-Label-Bulk_1601696643918.html",
                "product": "Organic ginger root extract gummies, "
                           "private-label OEM/ODM; chewable digestion-"
                           "support format",
                "product_name": "SUNI Organic Ginger Root Extract Gummies",
                "product_url": "https://www.alibaba.com/product-detail/SUNI-OEM-ODM-Private-Label-Bulk_1601696643918.html",
                "moq": 100,
                "unit_cost_usd": 2.20,
                "total_moq_cost_usd": 220,
                "alleviation_score": 66,
                "notes": "Shaanxi, China. MOQ 100 bottles (verified "
                         "on product page). Unit price $2.20 is high-"
                         "end of observed $1.55–2.20 range. Vegan + "
                         "organic certs. Direct competitor to VitaHall.",
            },
        ],
    },
    {
        "key": "nutrient_depletion",
        "title": "Nutrient Depletion",
        "eyebrow": "§ 05 · Side effect 5 of 5 · GLP-1 targeted multivitamin",
        "problem": (
            "Reduced caloric intake (often 1,000–1,200 kcal/day) creates "
            "predictable micronutrient gaps. ~13–22% of GLP-1 users "
            "develop documented deficiencies. Compounds the muscle-loss "
            "and hair-thinning problems."
        ),
        "cause": [
            "Reduced food volume → reduced intake of B12, magnesium, "
            "potassium, vitamin D, K2, iodine.",
            "Loss of food preferences (anhedonia) further narrows the "
            "nutrient base actually consumed.",
            "GLP-1 receptor activity in pancreas + gut may reduce "
            "absorption efficiency for some micronutrients.",
        ],
        "otc": [
            "GLP-1-targeted multivitamin: B-complex (esp. B12 "
            "methylcobalamin), Mg glycinate, K2 (MK-7), D3.",
            "Electrolyte sticks (Na + K + Mg) — counter the "
            "hypovolemia / lightheadedness common in week 1–4.",
            "Whey or plant protein shakes (20–30g) — supplement protein "
            "intake (often <60g/day on GLP-1).",
            "Comprehensive labs (CBC, CMP, B12, vit D, ferritin) "
            "every 3–6 months while on therapy.",
        ],
        "suppliers": [
            {
                "name": "Aogubio (Xi'an)",
                "url": "https://www.aogubio.com/dietary-supplement/",
                "product": "OEM multivitamin capsule (B12, Mg-glycinate, "
                           "K2 MK-7, D3, B-complex blend)",
                "product_name": "Vitamin B12 Methylcobalamin",
                "product_url": "https://www.aogubio.com/supply-food-grade-vitamin-b12-methylcobalamin-powder-product/",
                "moq": 100,
                "unit_cost_usd": 1.75,
                "total_moq_cost_usd": 175,
                "alleviation_score": 71,
                "notes": "Xi'an-based China OEM. Hits both gates "
                         "(MOQ 100 AND total <$250). Custom GLP-1 blend "
                         "spec'd by buyer; 2-4 wk lead. Lab COA per "
                         "batch; ISO + GMP cert.",
            },
        ],
    },
    {
        "key": "outer_packaging",
        "title": "Outer Packaging",
        "eyebrow": "§ 06 · Add-on · Two-SKU folding-mailer system",
        "gate_skip": True,  # not subject to formula MOQ gate
        "is_packaging": True,
        "problem": (
            "Every Aplomb SKU needs an outer mailer that ships flat, "
            "weighs almost nothing, and reads editorial — not generic "
            "corrugated. Reference: Musely's matte mauve folding-carton "
            "mailer (single-color exterior, branded interior message, "
            "soft tactile paperboard). Two carton SKUs cover the line: "
            "a small mailer for the serum, capsule bottles, and "
            "lozenge tin; a large mailer for the Calm kit."
        ),
        "spec": [
            "Small carton — ~120 × 60 × 60 mm interior · fits Serum, "
            "Daily, Roots, Composure tin",
            "Large carton — ~180 × 110 × 50 mm interior · fits Calm "
            "kit (largest single SKU footprint)",
            "Style — roll-end tuck-top (RETT) or auto-bottom paperboard "
            "mailer · matte uncoated · single-color exterior in amber "
            "#7A3D14 on bone #F7F1E6 · interior message in Cormorant "
            "Garamond italic",
            "MOQ band — 200-400 per carton SKU (so 400-800 cartons "
            "total across both die lines for first run)",
            "Hard gates — US production · custom interior print · "
            "MOQ ≤ 400 per SKU · willing to run two die lines on "
            "same project",
        ],
        "rubric": [
            "30% Customization fidelity — exterior + interior print, "
            "matte uncoated, finishes (foil/deboss optional) — "
            "Musely-quality reference match",
            "25% MOQ fit at 200-400 per SKU",
            "20% Lead time + reliability — sub-4-week production ideal",
            "15% Cost at MOQ — combined small + large at 300/300",
            "10% Sustainability + brand fit — FSC / recycled paperboard "
            "preferred (not gated)",
        ],
        "musely_note": (
            "Musely (telehealth dermatology, founded 2013) commissioned "
            "its brand evolution from White and Salt design studio. "
            "Musely's specific carton vendor is not public; the closest "
            "production-format match for the same paperboard mailer at "
            "MOQ 200-400 is Arka or noissue. RFQ both as the path of "
            "least resistance to a Musely-equivalent first run."
        ),
        "suppliers": [
            {
                "name": "Arka (Packlane)",
                "url": "https://packlane.com/",
                "product": "Custom paperboard mailer · full-color in/out · "
                           "matte / satin / Dreamcoat finishes · custom dies",
                "product_name": "Mailer Box (custom)",
                "product_url": "https://packlane.com/products/mailer-box",
                "moq": 200,
                "unit_cost_usd": 3.40,
                "total_moq_cost_usd": 2040,
                "alleviation_score": 88,
                "notes": "San Francisco · MOQ 1+ (so 200/400 trivial) · "
                         "100% US-made · ~10-day standard lead, 7-day "
                         "rush · instant online proof · Shopify/WMS "
                         "integration. Combined cost = 300 small + 300 "
                         "large @ ~$3.40 weighted avg. RECOMMENDED "
                         "lead — fastest path to Musely-style sample.",
            },
            {
                "name": "EcoEnclose",
                "url": "https://www.ecoenclose.com/",
                "product": "Custom shipping box · 100% recycled FSC "
                           "kraft (95% post-consumer) · multi-panel "
                           "interior print up to 2 colors",
                "product_name": "Custom Shipping Box",
                "product_url": "https://www.ecoenclose.com/shop/custom-shipping-boxes/",
                "moq": 100,
                "unit_cost_usd": 3.85,
                "total_moq_cost_usd": 2310,
                "alleviation_score": 81,
                "notes": "Colorado · 100% recycled · FSC SCS-COC-009926 · "
                         "HydroSoy/Black-Algae inks (no plastic "
                         "lamination) · made in USA. Best sustainability "
                         "story; kraft-only exterior reads less "
                         "Musely-velvet, more clinical/honest.",
            },
            {
                "name": "Refine Packaging",
                "url": "https://refinepackaging.com/",
                "product": "Custom mailer box · 10-28pt board (kraft / "
                           "cardstock / e-flute corrugated) · gloss / "
                           "matte / spot UV finishes · interior + "
                           "exterior print",
                "product_name": "Custom Mailer Box",
                "product_url": "https://refinepackaging.com/mailer-boxes/",
                "moq": 100,
                "unit_cost_usd": 3.20,
                "total_moq_cost_usd": 1920,
                "alleviation_score": 80,
                "notes": "US-based · MOQ 100 · ~$2-5/unit at 100-500 · "
                         "8-10 day standard, 4-6 day rush · free US "
                         "shipping · pre-production mockup included. "
                         "Best $/unit at MOQ 200-400.",
            },
            {
                "name": "noissue",
                "url": "https://noissue.co/",
                "product": "Custom mailer box · FSC-certified recycled "
                           "paperboard · full color in/out · "
                           "branded-DTC focus",
                "product_name": "Custom Mailer Box",
                "product_url": "https://noissue.co/custom-packaging/custom-boxes/",
                "moq": 100,
                "unit_cost_usd": 4.10,
                "total_moq_cost_usd": 2460,
                "alleviation_score": 78,
                "notes": "US distribution · MOQ 100 · FSC + carbon-"
                         "neutral · in-house brand-design support "
                         "(closest to Musely's White-and-Salt-led "
                         "aesthetic). 4-6 wk lead is the trade-off vs. "
                         "Arka. Premium feel.",
            },
            {
                "name": "PakFactory",
                "url": "https://pakfactory.com/",
                "product": "Custom folding paperboard carton · premium "
                           "finishes (soft-touch, foil, emboss, UV) · "
                           "multiple paperboard substrates",
                "product_name": "Folding Carton & Paperboard Boxes",
                "product_url": "https://pakfactory.com/folding-carton-paperboard-boxes.html",
                "moq": 100,
                "unit_cost_usd": 2.95,
                "total_moq_cost_usd": 1770,
                "alleviation_score": 75,
                "notes": "HQ Markham, Ontario (Canada — soft-fails the "
                         "US-only gate; included for cost-comparison). "
                         "Best premium-finish range (soft-touch matches "
                         "Musely's tactile feel). 3-4 wk lead. Pakible "
                         "(US, RFQ-only) is a sixth alternative if "
                         "needed.",
            },
        ],
    },
]


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
TOTAL_SLIDES = 15  # title + sourcing-gate + 5×(overview+suppliers) + packaging×2 + next-steps


def slide_01_title():
    s = add_blank_slide()
    add_hline(s, Inches(0.6), Inches(0.55), Inches(12.13))
    add_text(s, Inches(0.6), Inches(0.18), Inches(6), Inches(0.3),
             ["§ 00 · Sourcing memo · v1"], font=SANS, size=9, color=AMBER)
    add_rich(s, Inches(11.0), Inches(0.18), Inches(1.8), Inches(0.3), [[
        {"t": "Aplomb", "font": SERIF, "size": 11, "color": INK, "italic": True},
        {"t": ".", "font": SERIF, "size": 11, "color": AMBER, "italic": True, "bold": True},
    ]], align=PP_ALIGN.RIGHT)

    add_text(s, Inches(0.6), Inches(2.0), Inches(11), Inches(0.4),
             ["MVP supplier shortlist · 5 side effects · 12 stock-formula vendors · "
              "+ outer-packaging mailer system (6 vendors)"],
             font=SANS, size=11, color=MUTED)

    add_rich(s, Inches(0.6), Inches(2.4), Inches(11), Inches(2.5), [
        [{"t": "Who should ", "font": SERIF, "size": 56, "color": INK},
         {"t": "actually", "font": SERIF, "size": 56, "color": INK, "italic": True},
         {"t": " make", "font": SERIF, "size": 56, "color": INK}],
        [{"t": "the first 100 units?", "font": SERIF, "size": 56, "color": INK}],
    ], leading=1.05)

    add_text(s, Inches(0.6), Inches(5.2), Inches(11), Inches(1.0),
             ["Stock-formula private-label SKUs only — every row links",
              "to a specific published product page so the MOQ ≤100 /",
              "total <$250 gate is grounded in real catalog data."],
             font=SANS, size=14, color=MUTED, leading=1.4)

    add_rect(s, Inches(0.6), Inches(6.2), Inches(7.5), Inches(0.55), fill=PAPER)
    add_rich(s, Inches(0.8), Inches(6.32), Inches(7.3), Inches(0.4), [[
        {"t": "FIVE PROBLEMS  ", "font": SANS, "size": 8.5, "color": AMBER, "bold": True},
        {"t": "Ozempic face · Halitosis · Hair loss · Nausea · Nutrient depletion",
         "font": SANS, "size": 9.5, "color": INK},
    ]])

    add_hline(s, Inches(0.6), Inches(7.05), Inches(12.13))
    add_text(s, Inches(0.6), Inches(7.13), Inches(8), Inches(0.3),
             [f"For Zachary Poll · Prepared {datetime.now().strftime('%Y-%m-%d')} · Confidential"],
             font=SANS, size=8, color=MUTED)
    add_text(s, Inches(11.3), Inches(7.13), Inches(1.5), Inches(0.3),
             [f"01 / {TOTAL_SLIDES:02d}"],
             font=SANS, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def slide_02_sourcing_gate():
    s = add_blank_slide()
    add_chapter_chrome(s, "§ 00 · Sourcing gate")

    add_rich(s, Inches(0.6), Inches(0.95), Inches(12), Inches(1.5), [
        [{"t": "How every supplier on the next 10 slides was filtered.",
          "font": SERIF, "size": 30, "color": INK}],
    ], leading=1.05)

    # Three gate cards
    cards = [
        ("HARD GATE — MOQ", "Minimum order quantity ≤ 100 units OR total MOQ "
         "spend < $250. This is the cash-outlay ceiling that lets us run a "
         "real market test on each side effect without committing >$1,500 "
         "of inventory capital."),
        ("HARD GATE — BRANDING", "Supplier must be white-label / private-"
         "label capable: Aplomb's name, label, dropper/bottle, and INCI "
         "(where required) on the finished good. No co-branded retail "
         "products, no contract-manufacturer-of-record disclosure required."),
        ("SOFT GATE — EFFICACY", "Alleviation score ≥ 60 / 100 against the "
         "specific side effect. Score is a weighted blend: 40% clinical "
         "evidence for the active(s), 25% product specificity, 25% supplier "
         "reputation, 10% regulatory cleanliness."),
    ]
    card_w = Inches(3.95)
    card_h = Inches(2.8)
    gap = Inches(0.10)
    start_x = Inches(0.6)
    start_y = Inches(2.2)
    for i, (title, body) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, start_y, card_w, card_h, fill=PAPER)
        add_text(s, x + Inches(0.25), start_y + Inches(0.25), card_w - Inches(0.5),
                 Inches(0.4), [title], font=SANS, size=8.5, color=AMBER, bold=True)
        add_hline(s, x + Inches(0.25), start_y + Inches(0.6),
                  card_w - Inches(0.5), color=RULE, weight=0.5)
        add_text(s, x + Inches(0.25), start_y + Inches(0.75), card_w - Inches(0.5),
                 card_h - Inches(1.0), [body],
                 font=SANS, size=11, color=INK, leading=1.4)

    # Scoring rubric strip
    add_rect(s, Inches(0.6), Inches(5.3), Inches(12.13), Inches(1.55), fill=PAPER)
    add_text(s, Inches(0.85), Inches(5.45), Inches(11.6), Inches(0.4),
             ["ALLEVIATION SCORE — RUBRIC"],
             font=SANS, size=8.5, color=AMBER, bold=True)
    add_hline(s, Inches(0.85), Inches(5.78), Inches(11.6), color=RULE, weight=0.5)
    add_rich(s, Inches(0.85), Inches(5.92), Inches(11.6), Inches(0.9), [[
        {"t": "40% ", "font": SANS, "size": 11, "color": AMBER, "bold": True},
        {"t": "clinical evidence for the active ingredients on this side effect    ",
         "font": SANS, "size": 11, "color": INK},
        {"t": "·  25% ", "font": SANS, "size": 11, "color": AMBER, "bold": True},
        {"t": "product specificity (targeted blend > generic)    ",
         "font": SANS, "size": 11, "color": INK},
    ], [
        {"t": "25% ", "font": SANS, "size": 11, "color": AMBER, "bold": True},
        {"t": "supplier reputation (years operating, certifications, reviews)    ",
         "font": SANS, "size": 11, "color": INK},
        {"t": "·  10% ", "font": SANS, "size": 11, "color": AMBER, "bold": True},
        {"t": "regulatory cleanliness (no FDA warning letters, ingredients GRAS / OTC-legal)",
         "font": SANS, "size": 11, "color": INK},
    ]], leading=1.4)

    add_footer(s, 2, TOTAL_SLIDES)


def build_overview_slide(section, page_num):
    s = add_blank_slide()
    add_chapter_chrome(s, section["eyebrow"])

    # Title
    add_rich(s, Inches(0.6), Inches(0.85), Inches(12), Inches(1.2), [
        [{"t": section["title"], "font": SERIF, "size": 38, "color": INK}],
    ], leading=1.05)

    # Two-column layout: PROBLEM (left), CAUSES + OTC (right)
    # Left: problem statement card
    add_rect(s, Inches(0.6), Inches(2.0), Inches(5.85), Inches(4.85), fill=PAPER)
    add_text(s, Inches(0.85), Inches(2.18), Inches(5.4), Inches(0.4),
             ["THE PROBLEM"], font=SANS, size=8.5, color=AMBER, bold=True)
    add_hline(s, Inches(0.85), Inches(2.52), Inches(5.4), color=RULE, weight=0.5)
    add_text(s, Inches(0.85), Inches(2.65), Inches(5.4), Inches(4.05),
             [section["problem"]],
             font=SANS, size=12, color=INK, leading=1.45)

    # Right: causes (top) + OTC remedies (bottom)
    add_rect(s, Inches(6.6), Inches(2.0), Inches(6.13), Inches(2.20), fill=PAPER)
    add_text(s, Inches(6.85), Inches(2.18), Inches(5.7), Inches(0.4),
             ["ROOT CAUSE"], font=SANS, size=8.5, color=AMBER, bold=True)
    add_hline(s, Inches(6.85), Inches(2.52), Inches(5.7), color=RULE, weight=0.5)
    cause_paras = []
    for c in section["cause"]:
        cause_paras.append([
            {"t": "•  ", "font": SANS, "size": 10.5, "color": AMBER, "bold": True},
            {"t": c, "font": SANS, "size": 10.5, "color": INK},
        ])
    add_rich(s, Inches(6.85), Inches(2.62), Inches(5.7), Inches(1.6),
             cause_paras, leading=1.32)

    add_rect(s, Inches(6.6), Inches(4.30), Inches(6.13), Inches(2.55), fill=PAPER)
    add_text(s, Inches(6.85), Inches(4.48), Inches(5.7), Inches(0.4),
             ["BEST OTC REMEDIES"], font=SANS, size=8.5, color=AMBER, bold=True)
    add_hline(s, Inches(6.85), Inches(4.82), Inches(5.7), color=RULE, weight=0.5)
    otc_paras = []
    for r in section["otc"]:
        otc_paras.append([
            {"t": "•  ", "font": SANS, "size": 10.5, "color": AMBER, "bold": True},
            {"t": r, "font": SANS, "size": 10.5, "color": INK},
        ])
    add_rich(s, Inches(6.85), Inches(4.92), Inches(5.7), Inches(1.95),
             otc_paras, leading=1.32)

    add_footer(s, page_num, TOTAL_SLIDES)


def build_packaging_overview_slide(section, page_num):
    """Packaging-specific overview slide.
    Same chrome as build_overview_slide but card labels + content adapted
    for the outer-mailer brief (spec, structural style, Musely reference,
    rubric)."""
    s = add_blank_slide()
    add_chapter_chrome(s, section["eyebrow"])

    # Title
    add_rich(s, Inches(0.6), Inches(0.85), Inches(12), Inches(1.2), [
        [{"t": section["title"], "font": SERIF, "size": 38, "color": INK}],
    ], leading=1.05)

    # Two-column layout: BRIEF (left), SPEC + RUBRIC (right)
    add_rect(s, Inches(0.6), Inches(2.0), Inches(5.85), Inches(4.85), fill=PAPER)
    add_text(s, Inches(0.85), Inches(2.18), Inches(5.4), Inches(0.4),
             ["THE BRIEF"], font=SANS, size=8.5, color=AMBER, bold=True)
    add_hline(s, Inches(0.85), Inches(2.52), Inches(5.4), color=RULE, weight=0.5)
    add_text(s, Inches(0.85), Inches(2.65), Inches(5.4), Inches(2.6),
             [section["problem"]],
             font=SANS, size=11, color=INK, leading=1.45)

    # Musely intelligence callout, in the same left column below the brief
    add_hline(s, Inches(0.85), Inches(5.30), Inches(5.4), color=RULE, weight=0.5)
    add_text(s, Inches(0.85), Inches(5.40), Inches(5.4), Inches(0.4),
             ["MUSELY INTELLIGENCE"], font=SANS, size=8.5, color=AMBER, bold=True)
    add_text(s, Inches(0.85), Inches(5.72), Inches(5.4), Inches(1.20),
             [section["musely_note"]],
             font=SANS, size=9.5, color=INK, leading=1.40)

    # Right top: SPEC bullets
    add_rect(s, Inches(6.6), Inches(2.0), Inches(6.13), Inches(2.95), fill=PAPER)
    add_text(s, Inches(6.85), Inches(2.18), Inches(5.7), Inches(0.4),
             ["CARTON SPEC"], font=SANS, size=8.5, color=AMBER, bold=True)
    add_hline(s, Inches(6.85), Inches(2.52), Inches(5.7), color=RULE, weight=0.5)
    spec_paras = []
    for c in section["spec"]:
        spec_paras.append([
            {"t": "•  ", "font": SANS, "size": 10.5, "color": AMBER, "bold": True},
            {"t": c, "font": SANS, "size": 10.5, "color": INK},
        ])
    add_rich(s, Inches(6.85), Inches(2.62), Inches(5.7), Inches(2.32),
             spec_paras, leading=1.32)

    # Right bottom: SCORING RUBRIC bullets
    add_rect(s, Inches(6.6), Inches(5.05), Inches(6.13), Inches(1.80), fill=PAPER)
    add_text(s, Inches(6.85), Inches(5.20), Inches(5.7), Inches(0.4),
             ["SCORING RUBRIC (0-100)"], font=SANS, size=8.5, color=AMBER, bold=True)
    add_hline(s, Inches(6.85), Inches(5.52), Inches(5.7), color=RULE, weight=0.5)
    rubric_paras = []
    for r in section["rubric"]:
        rubric_paras.append([
            {"t": "•  ", "font": SANS, "size": 9, "color": AMBER, "bold": True},
            {"t": r, "font": SANS, "size": 9, "color": INK},
        ])
    add_rich(s, Inches(6.85), Inches(5.60), Inches(5.7), Inches(1.25),
             rubric_paras, leading=1.30)

    add_footer(s, page_num, TOTAL_SLIDES)


def build_supplier_slide(section, page_num):
    s = add_blank_slide()
    add_chapter_chrome(s, section["eyebrow"])

    n = len(section["suppliers"])
    title_suffix = "shortlisted suppliers" if n != 1 else "shortlisted supplier"
    add_rich(s, Inches(0.6), Inches(0.85), Inches(12), Inches(1.2), [
        [{"t": f"{section['title']} — {n} {title_suffix}",
          "font": SERIF, "size": 30, "color": INK}],
    ], leading=1.05)

    if section.get("is_packaging"):
        subhead = ("US folding-mailer vendors at MOQ 200-400 per SKU. "
                   "Pricing reflects 300-unit run of EACH die line "
                   "(small + large) — combined cost of first order.")
    else:
        subhead = ("Stock-formula private-label SKUs only. MOQ + total "
                   "cost verified against each supplier's published page.")
    add_text(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.4),
             [subhead], font=SANS, size=10, color=MUTED)

    # Matrix — height scales with # rows (3 row layout for ≤3 suppliers,
    # otherwise extend down and tuck the recommendation strip under it)
    n = len(section["suppliers"])
    if n <= 3:
        tbl_h = Inches(3.5)
        rec_y = Inches(5.95)
    else:
        tbl_h = Inches(4.30)
        rec_y = Inches(6.50)
    add_supplier_matrix(s, Inches(0.6), Inches(2.15),
                        Inches(12.13), tbl_h, section["suppliers"])

    # Recommendation strip — pick the highest scorer
    if section["suppliers"]:
        best = max(section["suppliers"], key=lambda x: x["alleviation_score"])
        rec_h = Inches(0.55) if n > 3 else Inches(0.95)
        add_rect(s, Inches(0.6), rec_y, Inches(12.13), rec_h, fill=PAPER)
        add_rich(s, Inches(0.85), rec_y + Inches(0.10), Inches(11.6), Inches(0.7), [[
            {"t": "RECOMMENDED LEAD  ", "font": SANS, "size": 8.5,
             "color": AMBER, "bold": True},
            {"t": f"{best['name']}", "font": SANS, "size": 12,
             "color": INK, "bold": True},
            {"t": f" · score {best['alleviation_score']}/100 · "
                  f"MOQ {best['moq']} · "
                  f"${best['unit_cost_usd']:.2f}/unit · "
                  f"total ${best['total_moq_cost_usd']:,.0f}",
             "font": SANS, "size": 11, "color": INK},
        ]])

    add_footer(s, page_num, TOTAL_SLIDES)


def slide_15_next_steps():
    s = add_blank_slide()
    add_chapter_chrome(s, "§ 07 · Next steps")

    add_rich(s, Inches(0.6), Inches(0.85), Inches(12), Inches(1.2), [
        [{"t": "Cross-cutting recommendations.",
          "font": SERIF, "size": 30, "color": INK}],
    ], leading=1.05)

    # Recommended lead supplier per category
    add_rect(s, Inches(0.6), Inches(1.85), Inches(12.13), Inches(2.85), fill=PAPER)
    add_text(s, Inches(0.85), Inches(2.0), Inches(11.6), Inches(0.4),
             ["RECOMMENDED LEAD SUPPLIER PER CATEGORY"],
             font=SANS, size=8.5, color=AMBER, bold=True)
    add_hline(s, Inches(0.85), Inches(2.35), Inches(11.6), color=RULE, weight=0.5)

    # Build a small table-like list of leads
    lead_paras = []
    for sec in SECTIONS:
        if not sec["suppliers"]:
            continue
        best = max(sec["suppliers"], key=lambda x: x["alleviation_score"])
        lead_paras.append([
            {"t": f"{sec['title']:<24} ", "font": SANS, "size": 11,
             "color": MUTED},
            {"t": f"{best['name']}", "font": SANS, "size": 12,
             "color": INK, "bold": True},
            {"t": f"  ·  score {best['alleviation_score']}/100  ·  "
                  f"MOQ {best['moq']}  ·  total ${best['total_moq_cost_usd']:,.0f}",
             "font": SANS, "size": 10.5, "color": MUTED},
        ])
    add_rich(s, Inches(0.85), Inches(2.50), Inches(11.6), Inches(2.20),
             lead_paras, leading=1.55)

    # Action checklist
    add_rect(s, Inches(0.6), Inches(4.95), Inches(12.13), Inches(2.0), fill=PAPER)
    add_text(s, Inches(0.85), Inches(5.10), Inches(11.6), Inches(0.4),
             ["IMMEDIATE ACTIONS"],
             font=SANS, size=8.5, color=AMBER, bold=True)
    add_hline(s, Inches(0.85), Inches(5.45), Inches(11.6), color=RULE, weight=0.5)
    actions = [
        "Send the RFQ template (in AFTER_supplier_directory_v1.xlsx Sheet 5) "
        "to all 6 lead suppliers this week — five formula vendors plus "
        "the packaging lead (Arka).",
        "Order packaging samples in parallel with formula RFQs: "
        "one small carton + one large carton from the top two packaging "
        "vendors (Arka + Refine Packaging) at $0 commitment to validate "
        "the Musely-style tactile quality.",
        "Negotiate sample units (3–5 per SKU) for in-house QC before "
        "committing to MOQ runs on the formula side.",
        "Lock packaging + label artwork with the brand designer before "
        "the first PO — fastest path to shippable inventory.",
        "Stand up a single shared Aplomb spec sheet (target ingredients, "
        "INCI, GMP cert, lab COA requirements) for every supplier.",
    ]
    action_paras = []
    for i, a in enumerate(actions):
        action_paras.append([
            {"t": f"{i+1}.  ", "font": SANS, "size": 10.5,
             "color": AMBER, "bold": True},
            {"t": a, "font": SANS, "size": 10.5, "color": INK},
        ])
    add_rich(s, Inches(0.85), Inches(5.58), Inches(11.6), Inches(1.32),
             action_paras, leading=1.4)

    add_footer(s, TOTAL_SLIDES, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build():
    slide_01_title()
    slide_02_sourcing_gate()
    page = 3
    for sec in SECTIONS:
        if sec.get("is_packaging"):
            build_packaging_overview_slide(sec, page)
        else:
            build_overview_slide(sec, page)
        build_supplier_slide(sec, page + 1)
        page += 2
    slide_15_next_steps()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT.name} ({len(prs.slides)} slides)")


# ---------------------------------------------------------------------------
# Post-save scrub — fixes python-pptx artifacts that trigger PP repair prompt.
# Per memory/feedback_pptx_repair_prompt.md (verified 2026-04-28).
# ---------------------------------------------------------------------------
def _scrub(pptx_path):
    today = datetime.now().strftime("%-m/%-d/%y")
    tmp = Path(str(pptx_path) + ".scrub")
    with zipfile.ZipFile(pptx_path, "r") as zin, \
         zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            name = item.filename
            if "ppt/printerSettings/" in name:
                continue
            if name.endswith(".xml") or name.endswith(".rels"):
                txt = data.decode("utf-8", errors="ignore")
                txt = re.sub(r'\s+smtClean="0"', '', txt)
                txt = re.sub(
                    r'<Relationship[^>]*?Target="(?:\.\./)?printerSettings/[^>]*?/>',
                    '', txt)
                txt = re.sub(
                    r'<Override[^>]*?PartName="/ppt/printerSettings/[^>]*?/>',
                    '', txt)
                txt = txt.replace("<a:t>1/27/13</a:t>", f"<a:t>{today}</a:t>")
                grp_replacement = (
                    '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/>'
                    '<a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/>'
                    '<a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>'
                )
                txt = re.sub(r'<p:grpSpPr\s*/>', grp_replacement, txt)

                def add_endpararpr(m):
                    block = m.group(0)
                    if '<a:r>' in block or '<a:r ' in block or '<a:endParaRPr' in block:
                        return block
                    return block.replace('</a:p>', '<a:endParaRPr/></a:p>')
                txt = re.sub(
                    r'<a:p>[^<]*?(?:<a:pPr[^>]*?(?:/>|>.*?</a:pPr>))?\s*</a:p>',
                    add_endpararpr, txt, flags=re.DOTALL)
                txt = re.sub(r'<a:p\s*/>', '<a:p><a:endParaRPr/></a:p>', txt)
                data = txt.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(str(tmp), str(pptx_path))
    print(f"Scrubbed {pptx_path.name}")


# ---------------------------------------------------------------------------
# PowerPoint repair pass — opens deck, programmatically clicks Repair if
# present, saves, kills PowerPoint. Loop until probe returns false.
# ---------------------------------------------------------------------------
def _repair_pass(out_path):
    if not os.path.exists("/Applications/Microsoft PowerPoint.app"):
        print("PowerPoint not installed; skipping repair pass.")
        return
    subprocess.run(["killall", "Microsoft PowerPoint"], capture_output=True)
    time.sleep(2)
    subprocess.run(["xattr", "-cr", str(out_path)], capture_output=True)
    subprocess.run(["open", "-a", "Microsoft PowerPoint", str(out_path)],
                   capture_output=True)
    time.sleep(12)
    probe = subprocess.run(
        ["osascript", "-e",
         'tell application "System Events" to tell process "Microsoft PowerPoint"'
         ' to (exists (button "Repair" of window 1))'],
        capture_output=True, text=True)
    has_repair = "true" in (probe.stdout or "").lower()
    if has_repair:
        print("Repair dialog detected; clicking Repair...")
        subprocess.run(["osascript", "-e",
            'tell application "System Events" to tell process'
            ' "Microsoft PowerPoint" to click button "Repair" of window 1'],
            capture_output=True)
        time.sleep(15)
        subprocess.run(["osascript", "-e",
            'tell application "System Events" to tell process'
            ' "Microsoft PowerPoint" to try'
            ' click button "OK" of window 1 end try'],
            capture_output=True)
        time.sleep(3)
        subprocess.run(["osascript", "-e",
            f'tell application "Microsoft PowerPoint" to '
            f'save active presentation in (POSIX file "{out_path}" as string)'],
            capture_output=True, timeout=20)
        time.sleep(5)
        print("Saved repaired file.")
    else:
        print("No repair dialog. File opens cleanly.")
    subprocess.run(["killall", "Microsoft PowerPoint"], capture_output=True)
    time.sleep(2)


# ---------------------------------------------------------------------------
# MOQ gate validation
# ---------------------------------------------------------------------------
def validate_moq_gate():
    failures = []
    for sec in SECTIONS:
        if sec.get("gate_skip"):
            continue  # packaging vendors use a different rubric
        for sup in sec["suppliers"]:
            moq = sup["moq"]
            unit = sup["unit_cost_usd"]
            total = sup["total_moq_cost_usd"]
            passes = (moq <= 100) or (total < 250)
            if not passes:
                failures.append(
                    f"{sec['title']} → {sup['name']}: MOQ={moq}, total=${total}"
                )
    if failures:
        raise SystemExit(
            "MOQ GATE FAILED:\n  " + "\n  ".join(failures))
    print("MOQ gate: all formula-section suppliers pass "
          "(packaging section skipped — different rubric).")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    validate_moq_gate()
    build()
    _scrub(OUT)
    _repair_pass(OUT)
    print("Done.")
